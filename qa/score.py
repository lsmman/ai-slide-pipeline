#!/usr/bin/env python3
"""QA scorer for an imagegen-ppt-pipeline deck.

This is the TEST (TDD): it defines the pass criteria for a finished deck and
scores the produced artifacts 0-100. Threshold to pass: 90.

Usage: python3 qa/score.py decks/<deck-name>
Exit 0 if score >= THRESHOLD else 1.
"""
import sys, os, json, glob

THRESHOLD = 90
CANVAS = (1920, 1080)


def load_slides(deck):
    with open(os.path.join(deck, "design", "slides.json"), encoding="utf-8") as f:
        return json.load(f)["slides"]


def png_size(path):
    from PIL import Image
    with Image.open(path) as im:
        return im.size


def pptx_slides(path):
    """Return list of (text_blob, notes_text) per slide."""
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for s in prs.slides:
        texts = []
        for sh in s.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
        notes = ""
        if s.has_notes_slide:
            notes = s.notes_slide.notes_text_frame.text or ""
        out.append((" \n".join(texts), notes))
    return out


def score(deck):
    slides = load_slides(deck)
    n = len(slides)
    report = {"deck": deck, "n_slides": n, "checks": {}, "issues": []}

    def add(key, got, total, issues=None):
        report["checks"][key] = {"got": round(got, 1), "max": total}
        if issues:
            report["issues"].extend(issues)

    # A. render completeness (20)
    a, ai = 0.0, []
    per = 20.0 / n
    for s in slides:
        p = os.path.join(deck, "render", f"{s['id']}.png")
        if os.path.exists(p):
            if png_size(p) == CANVAS:
                a += per
            else:
                ai.append(f"A {s['id']}: wrong size {png_size(p)}")
        else:
            ai.append(f"A {s['id']}: missing png")
    add("A_render_complete", a, 20, ai)

    # B. no overflow (20) — from render manifest
    b, bi = 0.0, []
    man = os.path.join(deck, "render", "manifest.json")
    if os.path.exists(man):
        m = json.load(open(man, encoding="utf-8"))
        mslides = {x["id"]: x for x in m.get("slides", [])}
        per = 20.0 / n
        for s in slides:
            ms = mslides.get(s["id"])
            if ms is None:
                bi.append(f"B {s['id']}: not in manifest")
            elif ms.get("overflow"):
                bad = [e["role"] for e in ms.get("elements", []) if e.get("overflow")]
                bi.append(f"B {s['id']}: overflow {bad}")
            else:
                b += per
    else:
        bi.append("B: render manifest missing")
    add("B_no_overflow", b, 20, bi)

    # C. image pptx (15): exists, n slides, notes present
    c, ci = 0.0, []
    ip = os.path.join(deck, "out", "deck-image.pptx")
    if os.path.exists(ip):
        ps = pptx_slides(ip)
        if len(ps) == n:
            c += 7.5
        else:
            ci.append(f"C image pptx slide count {len(ps)} != {n}")
        with_notes = sum(1 for _, nt in ps if len(nt.strip()) >= 100)
        c += 7.5 * (with_notes / n)
        if with_notes < n:
            ci.append(f"C image pptx notes: {with_notes}/{n} have >=100 chars")
    else:
        ci.append("C: deck-image.pptx missing")
    add("C_image_pptx", c, 15, ci)

    # D. editable pptx (20): exists, n slides, notes, title text present
    d, di = 0.0, []
    ep = os.path.join(deck, "out", "deck-editable.pptx")
    if os.path.exists(ep):
        ps = pptx_slides(ep)
        if len(ps) == n:
            d += 5
        else:
            di.append(f"D editable slide count {len(ps)} != {n}")
        with_notes = sum(1 for _, nt in ps if len(nt.strip()) >= 100)
        d += 5 * (with_notes / n)
        if with_notes < n:
            di.append(f"D editable notes: {with_notes}/{n}")
        # title text fidelity — check the rendered title-role element text
        # (slide.title is a short metadata label; the headline lives in elements)
        ok = 0
        for i, s in enumerate(slides):
            if i < len(ps):
                blob = ps[i][0].replace(" ", "").replace("\n", "")
                title_els = [e["text"] for e in s.get("elements", []) if e.get("role") == "title"]
                want = (title_els[0] if title_els else (s.get("title") or "")).replace(" ", "").replace("\n", "")
                if want and want in blob:
                    ok += 1
                else:
                    di.append(f"D {s['id']}: title element text not found in editable text")
        d += 10 * (ok / n)
    else:
        di.append("D: deck-editable.pptx missing")
    add("D_editable_pptx", d, 20, di)

    # E. pdf (10): exists + page count
    e, ei = 0.0, []
    pdf = os.path.join(deck, "out", "deck.pdf")
    pm = os.path.join(deck, "out", "pdf_manifest.json")
    if os.path.exists(pdf):
        e += 5
        if os.path.exists(pm):
            pages = json.load(open(pm, encoding="utf-8")).get("pages")
            if pages == n:
                e += 5
            else:
                ei.append(f"E pdf pages {pages} != {n}")
        else:
            ei.append("E: pdf_manifest.json missing")
    else:
        ei.append("E: deck.pdf missing")
    add("E_pdf", e, 10, ei)

    # F. explain docs (15): n md, each has title + region table
    f, fi = 0.0, []
    per = 15.0 / n
    for s in slides:
        p = os.path.join(deck, "explain", f"{s['id']}.md")
        if os.path.exists(p):
            txt = open(p, encoding="utf-8").read()
            title = s.get("title") or ""
            if title and title in txt and "| region" in txt:
                f += per
            else:
                fi.append(f"F {s['id']}: missing title or region table")
        else:
            fi.append(f"F {s['id']}: explain md missing")
    add("F_explain", f, 15, fi)

    total = a + b + c + d + e + f
    report["score"] = round(total, 1)
    report["threshold"] = THRESHOLD
    report["pass"] = total >= THRESHOLD
    return report


if __name__ == "__main__":
    deck = sys.argv[1] if len(sys.argv) > 1 else "decks/steelops-capstone"
    rep = score(deck)
    print(json.dumps(rep, ensure_ascii=False, indent=2))
    summary = " ".join(f"{k}={v['got']}/{v['max']}" for k, v in rep["checks"].items())
    print(f"\nSCORE {rep['score']}/100 (need {THRESHOLD})  {'PASS' if rep['pass'] else 'FAIL'}")
    print(summary)
    sys.exit(0 if rep["pass"] else 1)
